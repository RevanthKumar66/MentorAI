import io
import json
import logging
from typing import List, Optional

logger = logging.getLogger("mentorai-os.services.chunking_service")

class ChunkingService:
    """Service to extract plain text from various file formats and chunk it recursively."""

    def __init__(self, chunk_size: int = 800, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = ["\n\n", "\n", " ", ""]

    def extract_text(self, content: bytes, file_extension: str, mime_type: str = "") -> str:
        """Extracts plain text from the file content bytes based on extension/mime-type."""
        ext = file_extension.lower().strip().replace(".", "")
        
        # 1. Plaintext/Code/Markdown formats
        text_extensions = {
            "txt", "md", "py", "js", "jsx", "ts", "tsx", 
            "java", "cpp", "c", "go", "rs", "sql", "r", 
            "rmd", "yaml", "yml", "xml", "csv", "json"
        }
        
        if ext in text_extensions or mime_type.startswith("text/") or mime_type == "application/json":
            try:
                content_str = content.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    content_str = content.decode("latin-1")
                except Exception as e:
                    logger.error(f"Failed to decode text file: {e}")
                    raise ValueError("Failed to decode text file content")

            # Custom parsing for Jupyter Notebooks
            if ext == "ipynb":
                try:
                    notebook = json.loads(content_str)
                    cells_text = []
                    for cell in notebook.get("cells", []):
                        cell_type = cell.get("cell_type", "")
                        source = cell.get("source", [])
                        if isinstance(source, list):
                            source = "".join(source)
                        if source.strip():
                            cells_text.append(f"[{cell_type.upper()} CELL]\n{source.strip()}")
                    return "\n\n".join(cells_text)
                except Exception as e:
                    logger.warning(f"Failed to parse ipynb JSON structure, falling back to raw text: {e}")
                    return content_str
            
            return content_str

        # 2. PDF Document Format
        elif ext == "pdf" or mime_type == "application/pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content))
                text_parts = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"[Page {i+1}]\n{page_text.strip()}")
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"Error parsing PDF: {e}", exc_info=True)
                raise ValueError(f"Failed to parse PDF document: {str(e)}")

        # 3. Microsoft Word Word (docx) Document
        elif ext == "docx" or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            try:
                import docx
                doc = docx.Document(io.BytesIO(content))
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            text_parts.append(" | ".join(row_text))
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"Error parsing Word Document (docx): {e}", exc_info=True)
                raise ValueError(f"Failed to parse Word Document: {str(e)}")

        # 4. Microsoft PowerPoint Presentation (pptx)
        elif ext == "pptx" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                text_parts = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"Error parsing Presentation (pptx): {e}", exc_info=True)
                raise ValueError(f"Failed to parse Presentation slide structure: {str(e)}")

        # 5. Microsoft Excel Spreadsheet (xlsx)
        elif ext == "xlsx" or mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                text_parts = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    sheet_text = []
                    for row in ws.iter_rows(values_only=True):
                        row_str = " | ".join([str(val).strip() for val in row if val is not None])
                        if row_str.strip():
                            sheet_text.append(row_str)
                    if sheet_text:
                        text_parts.append(f"[Sheet {sheet}]\n" + "\n".join(sheet_text))
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"Error parsing Spreadsheet (xlsx): {e}", exc_info=True)
                raise ValueError(f"Failed to parse Spreadsheet content: {str(e)}")

        # 6. Fallback or unsupported formats
        else:
            logger.warning(f"Unparseable file type: .{ext} / {mime_type}, doing raw decode attempt")
            try:
                return content.decode("utf-8", errors="ignore")
            except Exception as e:
                raise ValueError(f"Unsupported file format: .{ext}")

    def split_text(self, text: str) -> List[str]:
        """Split text into chunks recursively based on separators."""
        return self._split_recursive(text, self.separators)

    def _split_recursive(self, text: str, separators: List[str]) -> List[str]:
        if len(text) <= self.chunk_size:
            return [text]

        # Find the first separator that exists in the text
        separator = separators[-1]
        new_separators = []
        for i, sep in enumerate(separators):
            if sep == "":
                separator = sep
                break
            if sep in text:
                separator = sep
                new_separators = separators[i+1:]
                break

        # Split text by separator
        if separator != "":
            splits = text.split(separator)
        else:
            splits = list(text)

        chunks = []
        current_chunk = []
        current_len = 0

        for split in splits:
            split_len = len(split) + (len(separator) if current_chunk else 0)
            
            if current_len + split_len > self.chunk_size:
                if current_chunk:
                    chunks.append(separator.join(current_chunk))
                    
                    # Backtrack to build overlap chunk
                    overlap_chunk = []
                    overlap_len = 0
                    for prev_split in reversed(current_chunk):
                        prev_len = len(prev_split) + (len(separator) if overlap_chunk else 0)
                        if overlap_len + prev_len > self.chunk_overlap:
                            break
                        overlap_chunk.insert(0, prev_split)
                        overlap_len += prev_len
                    
                    current_chunk = overlap_chunk
                    current_len = overlap_len
                
                # Handle single split larger than chunk_size
                if split_len > self.chunk_size:
                    sub_chunks = self._split_recursive(split, new_separators)
                    for sc in sub_chunks:
                        chunks.append(sc)
                    current_chunk = []
                    current_len = 0
                else:
                    current_chunk.append(split)
                    current_len += len(split)
            else:
                current_chunk.append(split)
                current_len += split_len

        if current_chunk:
            chunks.append(separator.join(current_chunk))

        return [c.strip() for c in chunks if c.strip()]
